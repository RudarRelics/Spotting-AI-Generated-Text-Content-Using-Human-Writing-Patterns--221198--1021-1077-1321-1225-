from flask import Flask, request, jsonify
from predict import predict_text
import os, io

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16 MB max upload

# ── Text extraction helpers ────────────────────────────────────────────────────

def extract_from_pdf(file_bytes):
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        return "\n".join(page.get_text() for page in doc).strip()
    except ImportError:
        raise RuntimeError("PyMuPDF not installed. Run: pip install pymupdf")

def extract_from_docx(file_bytes):
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except ImportError:
        raise RuntimeError("python-docx not installed. Run: pip install python-docx")

def extract_from_pptx(file_bytes):
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
        return "\n".join(texts)
    except ImportError:
        raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")

def extract_from_txt(file_bytes):
    return file_bytes.decode("utf-8", errors="ignore").strip()

EXTRACTORS = {
    ".pdf":  extract_from_pdf,
    ".docx": extract_from_docx,
    ".doc":  extract_from_docx,
    ".pptx": extract_from_pptx,
    ".ppt":  extract_from_pptx,
    ".txt":  extract_from_txt,
    ".md":   extract_from_txt,
    ".rtf":  extract_from_txt,
}

# ── Routes ─────────────────────────────────────────────────────────────────────

@app.route("/")
def home():
    return """
    <!DOCTYPE html>
    <html lang="en">
    <head>
        <meta charset="UTF-8">
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <title>AI Text Detector</title>
        <link href="https://fonts.googleapis.com/css2?family=Space+Mono:wght@400;700&family=Syne:wght@400;600;800&display=swap" rel="stylesheet">
        <style>
            *, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }

            :root {
                --bg: #0a0a0f;
                --surface: #12121a;
                --border: #1e1e2e;
                --accent: #00f5c4;
                --accent-dim: rgba(0,245,196,0.08);
                --accent-glow: rgba(0,245,196,0.25);
                --warn: #ff4f6d;
                --warn-dim: rgba(255,79,109,0.08);
                --warn-glow: rgba(255,79,109,0.25);
                --text: #e8e8f0;
                --muted: #5a5a7a;
                --font-mono: 'Space Mono', monospace;
                --font-sans: 'Syne', sans-serif;
            }

            body {
                background: var(--bg); color: var(--text);
                font-family: var(--font-sans); min-height: 100vh;
                display: flex; flex-direction: column; align-items: center;
                padding: 60px 20px; position: relative; overflow-x: hidden;
            }

            body::before {
                content: ''; position: fixed; inset: 0;
                background-image:
                    linear-gradient(rgba(0,245,196,0.03) 1px, transparent 1px),
                    linear-gradient(90deg, rgba(0,245,196,0.03) 1px, transparent 1px);
                background-size: 40px 40px; pointer-events: none; z-index: 0;
            }

            body::after {
                content: ''; position: fixed; top: -200px; left: 50%;
                transform: translateX(-50%); width: 800px; height: 600px;
                background: radial-gradient(ellipse, rgba(0,245,196,0.06) 0%, transparent 70%);
                pointer-events: none; z-index: 0;
            }

            .container { width: 100%; max-width: 760px; position: relative; z-index: 1; }

            header { text-align: center; margin-bottom: 50px; animation: fadeDown 0.7s ease both; }

            .badge {
                display: inline-block; font-family: var(--font-mono);
                font-size: 11px; letter-spacing: 3px; text-transform: uppercase;
                color: var(--accent); border: 1px solid var(--accent);
                padding: 5px 14px; border-radius: 2px; margin-bottom: 20px;
                background: var(--accent-dim);
            }

            h1 { font-size: clamp(2rem,5vw,3.2rem); font-weight: 800; line-height: 1.1; letter-spacing: -1px; color: #fff; }
            h1 span { color: var(--accent); }
            .subtitle { margin-top: 12px; color: var(--muted); font-size: 15px; font-family: var(--font-mono); }

            .card {
                background: var(--surface); border: 1px solid var(--border);
                border-radius: 12px; padding: 32px;
                animation: fadeUp 0.7s 0.15s ease both;
                position: relative; overflow: hidden;
            }
            .card::before {
                content: ''; position: absolute; top: 0; left: 0; right: 0; height: 1px;
                background: linear-gradient(90deg, transparent, var(--accent), transparent); opacity: 0.5;
            }

            /* Tabs */
            .tabs {
                display: flex; gap: 4px;
                background: var(--bg); border: 1px solid var(--border);
                border-radius: 8px; padding: 4px; margin-bottom: 24px;
            }
            .tab-btn {
                flex: 1; padding: 10px 0; background: none; border: none;
                border-radius: 6px; font-family: var(--font-mono); font-size: 12px;
                letter-spacing: 1px; text-transform: uppercase; color: var(--muted);
                cursor: pointer; transition: background 0.2s, color 0.2s;
                display: flex; align-items: center; justify-content: center; gap: 7px;
            }
            .tab-btn.active { background: var(--accent-dim); color: var(--accent); border: 1px solid rgba(0,245,196,0.2); }
            .tab-btn:not(.active):hover { color: var(--text); background: rgba(255,255,255,0.03); }

            .panel { display: none; }
            .panel.active { display: block; }

            .label {
                display: flex; align-items: center; gap: 8px;
                font-family: var(--font-mono); font-size: 11px;
                letter-spacing: 2px; text-transform: uppercase;
                color: var(--muted); margin-bottom: 12px;
            }
            .label::before {
                content: ''; width: 6px; height: 6px; border-radius: 50%;
                background: var(--accent); box-shadow: 0 0 8px var(--accent);
            }

            textarea {
                width: 100%; min-height: 220px;
                background: var(--bg); border: 1px solid var(--border);
                border-radius: 8px; color: var(--text);
                font-family: var(--font-mono); font-size: 14px; line-height: 1.7;
                padding: 18px; resize: vertical; outline: none;
                transition: border-color 0.2s, box-shadow 0.2s; caret-color: var(--accent);
            }
            textarea::placeholder { color: var(--muted); }
            textarea:focus { border-color: var(--accent); box-shadow: 0 0 0 3px var(--accent-dim); }

            .meta-row { display: flex; justify-content: space-between; align-items: center; margin-top: 10px; margin-bottom: 24px; }
            #charCount { font-family: var(--font-mono); font-size: 12px; color: var(--muted); }
            .clear-btn {
                font-family: var(--font-mono); font-size: 11px; color: var(--muted);
                background: none; border: none; cursor: pointer;
                letter-spacing: 1px; text-transform: uppercase; transition: color 0.2s; padding: 0;
            }
            .clear-btn:hover { color: var(--warn); }

            /* Drop zone */
            .drop-zone {
                border: 2px dashed var(--border); border-radius: 10px;
                padding: 48px 24px; text-align: center; cursor: pointer;
                transition: border-color 0.2s, background 0.2s;
                background: var(--bg); position: relative; margin-bottom: 16px;
            }
            .drop-zone:hover, .drop-zone.drag-over { border-color: var(--accent); background: var(--accent-dim); }
            .drop-zone input[type=file] {
                position: absolute; inset: 0; opacity: 0; cursor: pointer; width: 100%; height: 100%;
            }
            .drop-icon { font-size: 40px; margin-bottom: 14px; display: block; }
            .drop-title { font-size: 17px; font-weight: 700; color: var(--text); margin-bottom: 6px; }
            .drop-sub { font-family: var(--font-mono); font-size: 12px; color: var(--muted); }
            .drop-formats { display: flex; flex-wrap: wrap; justify-content: center; gap: 6px; margin-top: 16px; }
            .fmt-tag {
                font-family: var(--font-mono); font-size: 10px; letter-spacing: 1px;
                padding: 3px 9px; border-radius: 3px;
                border: 1px solid var(--border); color: var(--muted); text-transform: uppercase;
            }

            /* File chip */
            #fileChip {
                display: none; align-items: center; gap: 10px;
                background: var(--accent-dim); border: 1px solid rgba(0,245,196,0.2);
                border-radius: 8px; padding: 12px 16px; margin-bottom: 16px;
            }
            #fileChip.show { display: flex; }
            .chip-icon { font-size: 22px; flex-shrink: 0; }
            .chip-info { flex: 1; min-width: 0; }
            .chip-name { font-family: var(--font-mono); font-size: 13px; color: var(--accent); white-space: nowrap; overflow: hidden; text-overflow: ellipsis; }
            .chip-size { font-family: var(--font-mono); font-size: 11px; color: var(--muted); margin-top: 2px; }
            .chip-remove { background: none; border: none; color: var(--muted); cursor: pointer; font-size: 18px; padding: 0; transition: color 0.2s; flex-shrink: 0; }
            .chip-remove:hover { color: var(--warn); }

            /* Extracted preview */
            #extractedPreview {
                display: none; background: var(--bg); border: 1px solid var(--border);
                border-radius: 8px; padding: 16px;
                font-family: var(--font-mono); font-size: 12px; color: var(--muted);
                line-height: 1.7; max-height: 140px; overflow-y: auto;
                margin-bottom: 16px; white-space: pre-wrap;
            }
            #extractedPreview.show { display: block; }
            #extractedPreview::-webkit-scrollbar { width: 4px; }
            #extractedPreview::-webkit-scrollbar-thumb { background: var(--border); border-radius: 2px; }

            /* Buttons */
            .detect-btn {
                width: 100%; padding: 16px; background: var(--accent); color: var(--bg);
                border: none; border-radius: 8px; font-family: var(--font-sans);
                font-size: 15px; font-weight: 700; letter-spacing: 1px; cursor: pointer;
                position: relative; overflow: hidden; transition: transform 0.15s, box-shadow 0.2s;
            }
            .detect-btn::after { content: ''; position: absolute; inset: 0; background: rgba(255,255,255,0.15); opacity: 0; transition: opacity 0.2s; }
            .detect-btn:hover { transform: translateY(-2px); box-shadow: 0 8px 30px var(--accent-glow); }
            .detect-btn:hover::after { opacity: 1; }
            .detect-btn:active { transform: translateY(0); }
            .detect-btn:disabled { opacity: 0.5; cursor: not-allowed; transform: none; box-shadow: none; }

            .btn-inner { display: flex; align-items: center; justify-content: center; gap: 10px; }
            .spinner {
                width: 16px; height: 16px; border: 2px solid rgba(0,0,0,0.3);
                border-top-color: var(--bg); border-radius: 50%;
                animation: spin 0.7s linear infinite; display: none;
            }

            /* Result */
            #result {
                margin-top: 24px; border-radius: 10px; overflow: hidden;
                opacity: 0; transform: translateY(10px);
                transition: opacity 0.4s, transform 0.4s; pointer-events: none;
            }
            #result.show { opacity: 1; transform: translateY(0); pointer-events: auto; }

            .result-inner { padding: 24px 28px; border: 1px solid var(--border); border-radius: 10px; }
            .result-inner.human { background: var(--accent-dim); border-color: var(--accent); }
            .result-inner.ai    { background: var(--warn-dim);   border-color: var(--warn); }

            .result-label { font-family: var(--font-mono); font-size: 11px; letter-spacing: 3px; text-transform: uppercase; color: var(--muted); margin-bottom: 8px; }
            .result-verdict { font-size: 2rem; font-weight: 800; letter-spacing: -0.5px; }
            .result-inner.human .result-verdict { color: var(--accent); }
            .result-inner.ai    .result-verdict { color: var(--warn); }
            .result-sub { font-family: var(--font-mono); font-size: 13px; color: var(--muted); margin-top: 4px; }

            .source-badge {
                display: inline-flex; align-items: center; gap: 5px;
                font-family: var(--font-mono); font-size: 10px; letter-spacing: 1px;
                text-transform: uppercase; padding: 3px 9px; border-radius: 3px;
                border: 1px solid var(--border); color: var(--muted); margin-top: 8px;
            }

            .confidence-bar-wrap { margin-top: 18px; }
            .bar-label { display: flex; justify-content: space-between; font-family: var(--font-mono); font-size: 11px; color: var(--muted); margin-bottom: 6px; }
            .bar-track { height: 6px; background: var(--border); border-radius: 99px; overflow: hidden; }
            .bar-fill { height: 100%; border-radius: 99px; width: 0%; transition: width 0.8s cubic-bezier(0.16,1,0.3,1); }
            .result-inner.human .bar-fill { background: var(--accent); }
            .result-inner.ai    .bar-fill { background: var(--warn); }

            .tip {
                display: flex; align-items: flex-start; gap: 10px;
                margin-top: 20px; padding: 14px 16px;
                background: rgba(255,255,255,0.03); border-radius: 6px;
                border: 1px solid var(--border); font-family: var(--font-mono);
                font-size: 12px; color: var(--muted); line-height: 1.6;
            }
            .tip-icon { flex-shrink: 0; font-size: 15px; }

            footer {
                margin-top: 40px; text-align: center;
                font-family: var(--font-mono); font-size: 11px;
                color: var(--muted); letter-spacing: 1px;
                animation: fadeUp 0.7s 0.3s ease both;
            }

            @keyframes fadeDown { from { opacity:0; transform:translateY(-20px); } to { opacity:1; transform:translateY(0); } }
            @keyframes fadeUp   { from { opacity:0; transform:translateY(20px);  } to { opacity:1; transform:translateY(0); } }
            @keyframes spin     { to { transform:rotate(360deg); } }

            .toast {
                position: fixed; bottom: 30px; left: 50%;
                transform: translateX(-50%) translateY(80px);
                background: var(--warn); color: #fff;
                font-family: var(--font-mono); font-size: 13px;
                padding: 12px 22px; border-radius: 6px;
                transition: transform 0.3s; z-index: 999; pointer-events: none;
            }
            .toast.show { transform: translateX(-50%) translateY(0); }
            .toast.success { background: var(--accent); color: var(--bg); }
        </style>
    </head>
    <body>
        <div class="container">
            <header>
                <div class="badge">AI Text Detection &nbsp;·&nbsp; Human Writing Patterns</div>
                <h1>AI Text <span>Detector</span></h1>
                <p class="subtitle">// detecting AI using human writing patterns</p>
            </header>

            <div class="card">

                <div class="tabs">
                    <button class="tab-btn active" id="tabText" onclick="switchTab('text')">
                        ✏️ &nbsp;Paste Text
                    </button>
                    <button class="tab-btn" id="tabFile" onclick="switchTab('file')">
                        📎 &nbsp;Upload File
                    </button>
                </div>

                <!-- Text panel -->
                <div class="panel active" id="panelText">
                    <div class="label">Input Text</div>
                    <textarea id="text" placeholder="Paste or type any text here. The model will analyze linguistic patterns to determine if it was written by a human or generated by AI..."></textarea>
                    <div class="meta-row">
                        <span id="charCount">0 characters</span>
                        <button class="clear-btn" onclick="clearText()">Clear ×</button>
                    </div>
                    <button class="detect-btn" id="detectBtnText" onclick="detectText()">
                        <span class="btn-inner">
                            <span class="spinner" id="spinnerText"></span>
                            <span id="btnTextLabel">Run Detection</span>
                        </span>
                    </button>
                </div>

                <!-- File panel -->
                <div class="panel" id="panelFile">
                    <div class="label">Upload Document</div>

                    <div class="drop-zone" id="dropZone">
                        <input type="file" id="fileInput"
                               accept=".pdf,.doc,.docx,.ppt,.pptx,.txt,.md,.rtf"
                               onchange="handleFileSelect(this.files[0])">
                        <span class="drop-icon">📂</span>
                        <div class="drop-title">Drop your file here</div>
                        <div class="drop-sub">or click to browse from your device</div>
                        <div class="drop-formats">
                            <span class="fmt-tag">PDF</span>
                            <span class="fmt-tag">DOCX</span>
                            <span class="fmt-tag">DOC</span>
                            <span class="fmt-tag">PPTX</span>
                            <span class="fmt-tag">PPT</span>
                            <span class="fmt-tag">TXT</span>
                            <span class="fmt-tag">MD</span>
                            <span class="fmt-tag">RTF</span>
                        </div>
                    </div>

                    <div id="fileChip">
                        <span class="chip-icon" id="chipIcon">📄</span>
                        <div class="chip-info">
                            <div class="chip-name" id="chipName"></div>
                            <div class="chip-size" id="chipSize"></div>
                        </div>
                        <button class="chip-remove" onclick="removeFile()" title="Remove">✕</button>
                    </div>

                    <div id="extractedPreview"></div>

                    <button class="detect-btn" id="detectBtnFile" onclick="detectFile()" disabled>
                        <span class="btn-inner">
                            <span class="spinner" id="spinnerFile"></span>
                            <span id="btnFileLabel">Upload &amp; Analyze</span>
                        </span>
                    </button>
                </div>

                <!-- Result -->
                <div id="result">
                    <div class="result-inner" id="resultInner">
                        <div class="result-label">Verdict</div>
                        <div class="result-verdict" id="verdict"></div>
                        <div class="result-sub" id="resultSub"></div>
                        <div id="sourceBadge" class="source-badge"></div>
                        <div class="confidence-bar-wrap">
                            <div class="bar-label">
                                <span>Confidence</span>
                                <span id="confValue"></span>
                            </div>
                            <div class="bar-track">
                                <div class="bar-fill" id="barFill"></div>
                            </div>
                        </div>
                        <div class="tip">
                            <span class="tip-icon">💡</span>
                            <span id="tipText"></span>
                        </div>
                    </div>
                </div>
            </div>

            <footer>BUILT WITH PYTHON &amp; FLASK &nbsp;·&nbsp; LOCAL INFERENCE &nbsp;·&nbsp; NO DATA STORED</footer>
        </div>

        <div class="toast" id="toast"></div>

        <script>
        function switchTab(tab) {
            document.querySelectorAll('.tab-btn').forEach(b => b.classList.remove('active'));
            document.querySelectorAll('.panel').forEach(p => p.classList.remove('active'));
            const cap = tab.charAt(0).toUpperCase() + tab.slice(1);
            document.getElementById('tab' + cap).classList.add('active');
            document.getElementById('panel' + cap).classList.add('active');
            document.getElementById('result').classList.remove('show');
        }

        const textarea = document.getElementById('text');
        textarea.addEventListener('input', () => {
            const n = textarea.value.length;
            document.getElementById('charCount').textContent = n.toLocaleString() + ' character' + (n !== 1 ? 's' : '');
        });

        function clearText() {
            textarea.value = '';
            document.getElementById('charCount').textContent = '0 characters';
            document.getElementById('result').classList.remove('show');
        }

        function showToast(msg, type='error') {
            const t = document.getElementById('toast');
            t.textContent = msg;
            t.className = 'toast show' + (type === 'success' ? ' success' : '');
            setTimeout(() => t.classList.remove('show'), 3200);
        }

        const FILE_ICONS = { pdf:'📕', doc:'📘', docx:'📘', ppt:'📙', pptx:'📙', txt:'📄', md:'📄', rtf:'📄' };
        let selectedFile = null;

        function handleFileSelect(file) {
            if (!file) return;
            const ext = file.name.split('.').pop().toLowerCase();
            const allowed = ['pdf','doc','docx','ppt','pptx','txt','md','rtf'];
            if (!allowed.includes(ext)) { showToast('Unsupported file type.'); return; }
            if (file.size > 16 * 1024 * 1024) { showToast('File too large (max 16 MB).'); return; }
            selectedFile = file;
            document.getElementById('chipIcon').textContent = FILE_ICONS[ext] || '📄';
            document.getElementById('chipName').textContent = file.name;
            document.getElementById('chipSize').textContent = formatBytes(file.size);
            document.getElementById('fileChip').classList.add('show');
            document.getElementById('detectBtnFile').disabled = false;
            document.getElementById('extractedPreview').classList.remove('show');
            document.getElementById('result').classList.remove('show');
        }

        function removeFile() {
            selectedFile = null;
            document.getElementById('fileInput').value = '';
            document.getElementById('fileChip').classList.remove('show');
            document.getElementById('detectBtnFile').disabled = true;
            document.getElementById('extractedPreview').classList.remove('show');
            document.getElementById('result').classList.remove('show');
        }

        function formatBytes(b) {
            if (b < 1024) return b + ' B';
            if (b < 1048576) return (b/1024).toFixed(1) + ' KB';
            return (b/1048576).toFixed(2) + ' MB';
        }

        const dz = document.getElementById('dropZone');
        dz.addEventListener('dragover',  e => { e.preventDefault(); dz.classList.add('drag-over'); });
        dz.addEventListener('dragleave', () => dz.classList.remove('drag-over'));
        dz.addEventListener('drop', e => {
            e.preventDefault(); dz.classList.remove('drag-over');
            const file = e.dataTransfer.files[0];
            if (file) handleFileSelect(file);
        });

        async function detectText() {
            const text = textarea.value.trim();
            if (!text) { showToast('Please enter some text first.'); return; }
            if (text.length < 20) { showToast('Text too short — add at least 20 characters.'); return; }
            setLoading('text', true);
            try {
                const res = await fetch('/detect', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({ text })
                });
                if (!res.ok) throw new Error();
                const data = await res.json();
                showResult(data.prediction, data.confidence, 'Pasted Text');
            } catch { showToast('Error: Could not reach the server.'); }
            finally { setLoading('text', false); }
        }

        async function detectFile() {
            if (!selectedFile) { showToast('Please select a file first.'); return; }
            setLoading('file', true);
            const formData = new FormData();
            formData.append('file', selectedFile);
            try {
                const res = await fetch('/upload', { method: 'POST', body: formData });
                if (!res.ok) throw new Error();
                const data = await res.json();
                if (data.error) { showToast('Error: ' + data.error); return; }
                const preview = document.getElementById('extractedPreview');
                preview.textContent = data.extracted_text.slice(0, 600) + (data.extracted_text.length > 600 ? '…' : '');
                preview.classList.add('show');
                showResult(data.prediction, data.confidence, selectedFile.name);
                showToast('File analyzed successfully!', 'success');
            } catch { showToast('Error: Could not process file.'); }
            finally { setLoading('file', false); }
        }

        function setLoading(which, on) {
            const isText = which === 'text';
            const btn     = document.getElementById(isText ? 'detectBtnText' : 'detectBtnFile');
            const spinner = document.getElementById(isText ? 'spinnerText'   : 'spinnerFile');
            const label   = document.getElementById(isText ? 'btnTextLabel'  : 'btnFileLabel');
            btn.disabled = on;
            spinner.style.display = on ? 'block' : 'none';
            label.textContent = on ? 'Analyzing…' : (isText ? 'Run Detection' : 'Upload & Analyze');
        }

        function showResult(prediction, confidence, source) {
            const isAI = prediction.toLowerCase().includes('ai') || prediction.toLowerCase().includes('generated');
            document.getElementById('resultInner').className = 'result-inner ' + (isAI ? 'ai' : 'human');
            document.getElementById('verdict').textContent   = isAI ? '⚠ AI Generated' : '✓ Human Written';
            document.getElementById('resultSub').textContent = 'Model returned: "' + prediction + '"';
            document.getElementById('sourceBadge').textContent = '📎 Source: ' + source;
            document.getElementById('confValue').textContent  = confidence + '%';
            document.getElementById('tipText').textContent    = isAI
                ? 'This text shows patterns typical of large language models — uniform sentence structure, lack of personal voice, or overly smooth phrasing.'
                : 'This text shows natural human writing characteristics — varied rhythm, personal phrasing, and authentic expression.';
            document.getElementById('result').classList.add('show');
            setTimeout(() => { document.getElementById('barFill').style.width = confidence + '%'; }, 50);
        }
        </script>
    </body>
    </html>
    """

@app.route("/detect", methods=["POST"])
def detect():
    text = request.json.get("text", "")
    prediction, confidence = predict_text(text)
    return jsonify({"prediction": prediction, "confidence": confidence})


@app.route("/upload", methods=["POST"])
def upload():
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400

    file = request.files["file"]
    filename = file.filename or ""
    ext = os.path.splitext(filename)[1].lower()

    if ext not in EXTRACTORS:
        return jsonify({"error": f"Unsupported file type: {ext}"}), 400

    file_bytes = file.read()

    try:
        text = EXTRACTORS[ext](file_bytes)
    except RuntimeError as e:
        return jsonify({"error": str(e)}), 500
    except Exception as e:
        return jsonify({"error": f"Could not extract text: {str(e)}"}), 500

    if not text or len(text.strip()) < 20:
        return jsonify({"error": "Could not extract enough text from this file."}), 422

    prediction, confidence = predict_text(text)
    return jsonify({
        "prediction": prediction,
        "confidence": confidence,
        "extracted_text": text,
        "char_count": len(text)
    })


if __name__ == "__main__":
    app.run(debug=True)